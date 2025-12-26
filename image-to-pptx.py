from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os
import re

def natural_sort_key(s):
    """Extract numbers from filename for natural sorting"""
    return [int(text) if text.isdigit() else text.lower()
            for text in re.split('([0-9]+)', s)]

def fit_image_to_slide(slide, image_path):
    """Add image to slide and fit it within the slide boundaries"""
    # Get actual slide dimensions from the presentation
    prs = slide.part.package.presentation_part.presentation
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Get image dimensions
    img = Image.open(image_path)
    img_width, img_height = img.size
    
    # Calculate aspect ratios
    img_aspect = img_width / img_height
    slide_aspect = slide_width / slide_height
    
    # Fit image within slide while maintaining aspect ratio
    if img_aspect > slide_aspect:
        # Image is wider - fit to width
        width = slide_width
        height = slide_width / img_aspect
    else:
        # Image is taller - fit to height
        height = slide_height
        width = slide_height * img_aspect
    
    # Center the image on the slide
    left = (slide_width - width) / 2
    top = (slide_height - height) / 2
    
    slide.shapes.add_picture(image_path, left, top, width=width, height=height)

def process_subdirectory(subdir_path, template_path, output_path):
    """Process a single subdirectory and create a PowerPoint"""
    
    # Load the template for this subdirectory
    prs = Presentation(template_path)
    
    # Get all image files
    image_extensions = ('.png', '.jpg', '.jpeg', '.bmp', '.gif', '.tiff')
    images = [f for f in os.listdir(subdir_path) 
              if f.lower().endswith(image_extensions)]
    images.sort(key=natural_sort_key)
    
    print(f"  Found {len(images)} images")
    
    slide_index = 0
    
    # Add images to slides
    for image in images:
        image_path = os.path.join(subdir_path, image)
        
        # Use existing slide or create new one
        if slide_index < len(prs.slides):
            slide = prs.slides[slide_index]
        else:
            # Add new blank slide
            blank_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_layout)
        
        # Add image to slide
        try:
            fit_image_to_slide(slide, image_path)
            print(f"  Added: {image} to slide {slide_index + 1}")
            slide_index += 1
        except Exception as e:
            print(f"  Error adding {image}: {e}")
    
    # Save the presentation
    prs.save(output_path)
    print(f"  Saved: {output_path} with {slide_index} slides\n")

def process_all_directories(base_dir, template_path, output_dir):
    """Process all subdirectories and create separate PowerPoints"""
    
    # Create output directory if it doesn't exist
    os.makedirs(output_dir, exist_ok=True)
    
    # Get all subdirectories
    subdirs = [d for d in os.listdir(base_dir) 
               if os.path.isdir(os.path.join(base_dir, d))]
    subdirs.sort(key=natural_sort_key)
    
    print(f"Found {len(subdirs)} subdirectories\n")
    
    # Process each subdirectory separately
    for subdir in subdirs:
        subdir_path = os.path.join(base_dir, subdir)
        output_path = os.path.join(output_dir, f"{subdir}.pptx")
        
        print(f"Processing: {subdir}")
        
        try:
            process_subdirectory(subdir_path, template_path, output_path)
        except Exception as e:
            print(f"  Error processing {subdir}: {e}\n")
    
    print(f"Complete! Created {len(subdirs)} PowerPoint files in: {output_dir}")

# Example usage
if __name__ == "__main__":
    # Configure these paths
    BASE_DIRECTORY = r"path"  # Directory containing subdirectories
    TEMPLATE_FILE = r"path"  # Your pre-built PowerPoint
    OUTPUT_DIRECTORY = r"path"  # Where to save the result
    
    # Run the script
    process_all_directories(BASE_DIRECTORY, TEMPLATE_FILE, OUTPUT_DIRECTORY)
